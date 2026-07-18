import io
from dataclasses import dataclass

from pptx import Presentation
from pypdf import PdfReader
from pypdf.errors import FileNotDecryptedError, PdfReadError

from app.errors import AppError


@dataclass(frozen=True)
class ParsedPage:
    page_no: int
    text: str
    ocr_used: bool = False


@dataclass(frozen=True)
class ParsedDeck:
    pages: list[ParsedPage]
    parser: str
    partial: bool = False

    @property
    def text(self) -> str:
        return "\n\n".join(page.text for page in self.pages if page.text)


def parse_pdf(data: bytes) -> ParsedDeck:
    try:
        reader = PdfReader(io.BytesIO(data))
        if reader.is_encrypted:
            raise AppError(
                "DOCUMENT_ENCRYPTED",
                "Encrypted PDFs are not supported.",
                status_code=415,
            )
        pages = [
            ParsedPage(index + 1, (page.extract_text() or "").strip())
            for index, page in enumerate(reader.pages)
        ]
        return ParsedDeck(pages=pages, parser="pypdf", partial=any(not page.text for page in pages))
    except AppError:
        raise
    except (PdfReadError, FileNotDecryptedError) as exc:
        raise AppError("DECK_UNPARSEABLE", "The PDF could not be parsed.", status_code=422) from exc


def parse_pptx(data: bytes) -> ParsedDeck:
    try:
        presentation = Presentation(io.BytesIO(data))
    except Exception as exc:
        raise AppError(
            "DECK_UNPARSEABLE", "The PPTX could not be parsed.", status_code=422
        ) from exc
    pages: list[ParsedPage] = []
    for index, slide in enumerate(presentation.slides, 1):
        lines = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text.strip():
                lines.append(shape.text.strip())
        pages.append(ParsedPage(index, "\n".join(lines)))
    return ParsedDeck(
        pages=pages, parser="python-pptx", partial=any(not page.text for page in pages)
    )


def parse_deck(data: bytes, mime: str) -> ParsedDeck:
    if mime == "application/pdf":
        return parse_pdf(data)
    if mime == "application/vnd.openxmlformats-officedocument.presentationml.presentation":
        return parse_pptx(data)
    raise AppError(
        "UNSUPPORTED_MEDIA_TYPE", "Only PDF and PPTX documents are accepted.", status_code=415
    )
