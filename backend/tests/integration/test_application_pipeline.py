import asyncio
import io

import pytest
from pptx import Presentation

from tests.conftest import signup_client


def sample_deck() -> bytes:
    presentation = Presentation()
    slides = [
        "Dolphin makes venture diligence evidence-first",
        "ARR $1.2M in Q4",
        "We serve 120 customers",
        "TAM $10B across global early-stage investing",
    ]
    for text in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[5])
        box = slide.shapes.add_textbox(0, 0, 5_000_000, 1_000_000)
        box.text = text
    output = io.BytesIO()
    presentation.save(output)
    return output.getvalue()


@pytest.mark.integration
async def test_application_to_evidence_backed_memo(client) -> None:
    auth = await signup_client(client)
    headers = {
        "Authorization": f"Bearer {auth['access_token']}",
        "Idempotency-Key": "application-golden-path",
    }
    response = await client.post(
        "/api/v1/applications",
        headers=headers,
        data={"company_name": "Dolphin Systems"},
        files={
            "deck": (
                "dolphin.pptx",
                sample_deck(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    )
    assert response.status_code == 202, response.text
    payload = response.json()

    job = None
    for _ in range(80):
        await asyncio.sleep(0.05)
        job_response = await client.get(
            f"/api/v1/jobs/{payload['job_id']}",
            headers={"Authorization": f"Bearer {auth['access_token']}"},
        )
        job = job_response.json()
        if job["status"] in {"succeeded", "failed", "partial"}:
            break
    assert job is not None
    assert job["status"] == "succeeded", job

    detail = await client.get(
        f"/api/v1/opportunities/{payload['opportunity_id']}",
        headers={"Authorization": f"Bearer {auth['access_token']}"},
    )
    assert detail.status_code == 200, detail.text
    body = detail.json()
    assert "overall_score" not in body
    assert set(body["axes"]) == {"founder", "market", "idea_vs_market"}
    assert all(body["axes"][axis] is not None for axis in body["axes"])
    assert body["memo"]["status"] == "ready"

    claims = await client.get(
        f"/api/v1/opportunities/{payload['opportunity_id']}/claims",
        headers={"Authorization": f"Bearer {auth['access_token']}"},
    )
    assert claims.status_code == 200
    assert len(claims.json()) >= 4
    for claim in claims.json():
        evidence = await client.get(
            f"/api/v1/opportunities/{payload['opportunity_id']}/evidence/{claim['id']}",
            headers={"Authorization": f"Bearer {auth['access_token']}"},
        )
        assert evidence.status_code == 200
        assert evidence.json(), claim


@pytest.mark.integration
async def test_idempotent_application_replay(client) -> None:
    auth = await signup_client(client)
    headers = {
        "Authorization": f"Bearer {auth['access_token']}",
        "Idempotency-Key": "repeat-me",
    }
    request = {
        "headers": headers,
        "data": {"company_name": "Replay Co"},
        "files": {
            "deck": (
                "deck.pptx",
                sample_deck(),
                "application/vnd.openxmlformats-officedocument.presentationml.presentation",
            )
        },
    }
    first = await client.post("/api/v1/applications", **request)
    second = await client.post("/api/v1/applications", **request)
    assert first.status_code == second.status_code == 202
    assert first.json() == second.json()
    assert second.headers["Idempotent-Replay"] == "true"
